"""
Chart Transfer
Loads data from a specially formatted excel workbook
and creates native Powerpoint charts from the Excel data in a new presentation.
Default file locations and chart template names stored in config.txt file.
A table in the "index" worksheet gives the sheet names for each set of chart data_name
The "sheets" col has the sheet names and the "type" col has each chart type
Other data can be added to this sheet over time to alter formatting of each chart.
Dependent on chart_type_data.py file which contains PPT chart type data.
"""
#import library and library files
from string import Template
#from numpy.testing._private.utils import print_assert_equal
import pandas as pd
from pptx import Presentation
from pptx.chart.data import CategoryChartData
#from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches
import tkinter as tk
from tkinter import *
from tkinter import filedialog
import sys
import os

# Import dictionary files that contain chart type data for ppt - these contain all the chart types PPT uses to make it universal
from chart_type_data import chart_type_dict, chart_template_dict


class filelocations:
    """class to import defaults and store file locations and values for individual chart templates
        uses a configuration file for the defaults.
    """
    def __init__(self):
        # import defaults
        self.import_defaults()
        
    def import_defaults(self):
        #import default values from config.txt and set up class attributes
        import configparser
        self.config = configparser.ConfigParser()
        self.config.read("config.txt")
        defaults = self.config['DEFAULTS']
        self.layout_number = int(defaults['layout_number'])
        self.chart_template_location = defaults['chart_template_location']
        self.working_folder = defaults['working_folder']
        self.chart_types=(list(dict(self.config["CHART TEMPLATES"].items())))
        
        self.chart_templates = list(dict(self.config['CHART TEMPLATES']).values())
        self.output_name=defaults["output_name"]  #r"C:\Users\jamie\OneDrive - Market Prescience\python\powerpoint\output.pptx"
        self.data_name=defaults["data_name"]  #"default_data.xlsx"
        self.template_name=defaults["template_name"]  #"default_template.pptx"

    def write_defaults(self):
        #export config setting to config.txt
        defaults = self.config['DEFAULTS']
        defaults['chart_template_location'] = self.chart_template_location
        defaults["output_name"] = self.output_name  
        defaults["data_name"] = self.data_name  
        defaults["template_name"] = self.template_name
        defaults['layout_number'] = str(self.layout_number)
        for template,type in zip(self.chart_templates,self.chart_types):
            self.config['CHART TEMPLATES'][type] = template 




        #write the default data to config file   
        with open ("config.txt", "w") as configfile:
            self.config.write(configfile)

class MainApp(tk.Tk):
    """ class for the main menu and application
    """

    def __init__(self):
        super().__init__() #inherit tk

        #Create file location data with defaults from file location class
        self.all_files = filelocations()

        # create main window and buttons
        self.title("Chart Transfer 1.0")
        self.geometry('400x50')
        self.show_defaults = False
        menu = Menu(self)
        self.config(menu=menu)
        self.rowconfigure([0,1,2,3,4,5,6,7,8,9,10,11,12,13,14], minsize=25)
        self.columnconfigure([0, 1], minsize=200)
        self.buttons = []

        #Make labels for top of cols
        self.make_label("Defaults",2,0,150,bg_col="skyblue")
        self.make_label("File locations (editable)",2,1,150,bg_col="skyblue")

        # Make the static buttons for editing default file locations
        self.make_button("Press here to start transfer",self.transfer_all,0,0,bg_col="red")
        self.make_button("Toggle show file names",self.show_files,0,1,bg_col="black")  
        self.make_button(self.all_files.template_name,self.choose_template,3,1)
        self.make_button(self.all_files.data_name,self.choose_data,4,1)
        self.make_button(self.all_files.output_name,self.choose_output,5,1)


        #Make labels for the buttons on the left
        self.make_label("Template presentation:",3,0,150)
        self.make_label("Excel data file:",4,0,150)
        self.make_label("Output presentation:",5,0,150)
        self.make_label("Layout number:",6,0,150)
        #Make menu for layout number
        self.clicked = IntVar()
        self.clicked.set(self.all_files.layout_number)        
        self.drop = OptionMenu(self,self.clicked,0,1,2,3,4,5,6,7,8,9,10)
        self.drop.grid(row=6, column=1, sticky="nsew")
        self.drop.config(bg="grey", fg="white")


        # Make the buttons for editing chart template file locations
        for index,chart_template in enumerate(self.all_files.chart_templates):
            self.button = tk.Button(self,text=chart_template, bg="dark slate grey",  fg="white", anchor="w", command = lambda j=index: self.edit_file(j))
            self.button.grid(row=index+7, column=1, sticky="nsew")
            self.buttons.append(self.button)
            #self.make_button(chart_type,self.edit_file,index,0)
        

        #Create file location data with defaults from file location class
        for index,chart_type in enumerate(self.all_files.chart_types):
            self.make_label(f"Template: {chart_type}",index+7,0,150, bg_col="dark slate grey")

        # create file menu
        file_menu = Menu(menu)
        menu.add_cascade(label='File', menu=file_menu)

        #file_menu.add_command(label='Options', command=options)
        file_menu.add_command(label='Save defaults',command=self.save_defaults)
        file_menu.add_command(label='Exit',command=self.exit_prog)


        # create run menu
        run_menu = Menu(menu)
        menu.add_cascade(label='Run', menu=run_menu)
        run_menu.add_command(label='Transfer', command=self.transfer_all)

        # create help menu
        help_menu = Menu(menu)
        menu.add_cascade(label='Help', menu=help_menu)
        help_menu.add_command(label='Help', command=self.show_help)
        help_menu.add_command(label="About", command=self.show_about)

    def show_files(self):
        if self.show_defaults==False:
            self.geometry('400x392')
            self.show_defaults = True
        else:
            self.geometry('400x50')
            self.show_defaults = False



    def set_layout(self):
        pass
    
    def edit_file(self,index):
        # open file dialog to select new template file
        self.all_files.chart_templates[index] = os.path.split(filedialog.askopenfilename(title="Select a File",initialdir=self.all_files.chart_template_location, filetypes=(("MS Chart template files","*.crtx"),("All types","*.*"))))[1]
        self.buttons[index+5].configure(text=self.all_files.chart_templates[index])

    def show_about(self):
        #display about text
        from tkinter import messagebox
        from help_text import about
        messagebox.showinfo("About",about)   

    def show_help(self):
        #display help text
        from tkinter import messagebox
        from help_text import help
        messagebox.showinfo("Help",help)    
   
    def save_defaults(self):
        self.all_files.layout_number=self.clicked.get() # reset the layout number if the default is saved w/o running the transfer
        self.all_files.write_defaults()

    def edit_defaults(self):
        self.template_window = ChangeTemplateDefaults(self.all_files.chart_templates,self.all_files.chart_types,self.all_files.chart_template_location)
        self.all_files.chart_templates = self.template_window.chart_templates

    def make_button(self,main_text,cmd,row,col,bg_col="grey",fg_col="white"):
        #Make a tk button and place in the correct location
        self.button = tk.Button(text=main_text, bg=bg_col, fg=fg_col, command = cmd,anchor="w")
        self.buttons.append(self.button)
        self.button.grid(row=row, column=col, sticky="nsew")

    def make_label(self,main_text,row,col,length,bg_col="grey",fg_col="white"):
        #make a tk label and place it in the correct location
        label = tk.Label(text=main_text, bg=bg_col, fg=fg_col, wraplength = length,relief="ridge",anchor="w")
        label.grid(row=row, column=col, sticky="nsew")

    # open file menu to select template file - just a plane pptx file rather than a pptm
    def choose_template(self):
        self.all_files.template_name = os.path.split(filedialog.askopenfilename(title="Select a File",initialdir=self.all_files.working_folder,filetypes=(("Powerpoint","*.pptx"),("All types","*.*"))))[1]
        self.buttons[2].configure(text=self.all_files.template_name)
        

    # open file menu to select data file
    def choose_data(self):
        self.all_files.data_name = os.path.split(filedialog.askopenfilename(title="Select a File",initialdir=self.all_files.working_folder,filetypes=(("Excel","*.xlsx"),("All types","*.*"))))[1]
        self.buttons[3].configure(text=self.all_files.data_name)
        

    # open file menu to select output pointpoint file - will be written over
    def choose_output(self):
        self.all_files.output_name = os.path.split(filedialog.askopenfilename(title="Select a File",initialdir=self.all_files.working_folder,filetypes=(("Powerpoint","*.pptx"),("All types","*.*"))))[1]
        self.buttons[4].configure(text=self.all_files.output_name)
        
       
    def exit_prog(self):
        # exit program 
        sys.exit()

    # tests whether there is a template, a data file and a output destination
    def transfer_all(self):
        if self.all_files.output_name=="" or self.all_files.data_name =="" or self.all_files.template_name =="":
            self.make_label("Missing file",1,0,150)
        else:
            self.execute_transfer()

    def open_chart_index(self,file_name):
        #open the excel file and grab the chart index infomation from index worksheet
        #this data contains the excel sheet names for each chart and chart types
        df = pd.read_excel (file_name, sheet_name="index",engine='openpyxl')
        return df

    def open_presentation(self,template):
        #open the presentation template - this will be user defined
        prs = Presentation(template)
        return prs

    def chart_data_load(self,file_name,sheet):
        # Open excel file and load data from sheet into pandas data frame
        df = pd.read_excel (file_name, sheet_name=sheet,engine='openpyxl')
        return df

    def create_chart(self,prs,df,col_names,chart_type,chart_number):
        # each chart is on indiviudal slide so create the slide and add a chart object than the chart data from excel
        self.all_files.layout_number=self.clicked.get()
        slide = prs.slides.add_slide(prs.slide_layouts[self.all_files.layout_number]) # select the layout - using 5 for default template but can be userdefined 
        chart_data = CategoryChartData()
        #Take the category col from the excel chart table and add to the active chart data as the .catgories 
        chart_data.categories = (df['category'].values.tolist())
        # take the category column off the chart data list
        col_names.pop(0)

        #count col and row Numbers
        num_of_cols = len(col_names)
        num_of_rows = len(df['category'].values.tolist())

        # loop through the columns of data and add to the active chart_data note that the data needs to be converted from list to tuple
        for x in range(num_of_cols):
            chart_data.add_series(str(col_names[x]),tuple(df[col_names[x]].values.tolist()))

        # add chart to slide --------------------
        x, y, cx, cy = Inches(1), Inches(2), Inches(10.5), Inches(5)
        # lookup chart_type info from dictionary of excel chart types
        obj = chart_type_dict[chart_type]
        chart=slide.shapes.add_chart(
            obj, x, y, cx, cy, chart_data
            ).chart
        slide.shapes.Name = "chart"+str(chart_number)
        return prs,chart,num_of_cols,num_of_rows

    def apply_chart_templates(self,types,labels,num_of_cols,num_of_rows):
        # Uses win32com to open powerpoint and apply chart templates to all the charts in the output file
        # Import the library and load ppt with the right file
        import win32com.client
        from pathlib import Path # import Path
        pptApp = win32com.client.Dispatch("Powerpoint.Application")
        
        pptApp.Visible = True
        print("output:",self.all_files.output_name)
        file = Path(self.all_files.working_folder+self.all_files.output_name) # use Path to make sure filename in correct format
        prs = pptApp.Presentations.Open(file, ReadOnly = False)

        # go through all the slides in the Presentation
        # examine each of the "Shapes" on slide - charts are in the Shapes object of a ppt slide 
        # find all the charts and then apply the template
        for index,current_slide in enumerate(prs.Slides):
            for shp in current_slide.Shapes:
                if shp.Type == 3:   #type 3 is msochart so if shape is a chart
                    # Apply the correct template to the chart
                    x=chart_template_dict[types[index-1]]
                    shp.Chart.ApplyChartTemplate(Path(self.all_files.chart_template_location+self.all_files.chart_templates[x-1]))
                    shp.Chart.Axes(2).MaximumScaleIsAuto = True # 2 - is the xlValue constant number so .Axes(2) is the x-axis
                    shp.Chart.Axes(2).MinimumScaleIsAuto = True
                  
                    #Loop through the columns in the data for the chart making sure the labels match the index sheet in the data file (numbers or percentages)
                    # Important as only one template type per chart type - could have separate chart templates for each variant but this is simipler(?)
                    for col in range(1,num_of_cols[index-1]+1):
                        plot = shp.Chart.SeriesCollection(col)
                        plot.HasDataLabels = True
                        for row in range(1,num_of_rows[index-1]+1):
                            if labels[index-1]=="per":
                                plot.DataLabels(row).NumberFormat = "0%"
                            else:
                                plot.DataLabels(row).NumberFormat = "0"

    def execute_transfer(self):
        #load the chart index dataframe
        chart_index_df = self.open_chart_index(self.all_files.data_name)
        
        #open the template presentation
        main_presentation = self.open_presentation(self.all_files.template_name)

        #create all the charts from the chart index
        try:
            chart_sheets = chart_index_df['sheets'].values.tolist() # put the sheet names into list
            chart_types = chart_index_df['type'].values.tolist() #put chart types into list
            chart_labels = chart_index_df['label'].values.tolist() # put the chart label type into list (numbers or counts are num and % are per)
            self.make_label("Data file valid",1,0,150)
        except:
            self.make_label("Data file not valid",1,0,150)
            return
        
        
        num_of_rows = []
        num_of_cols = []
 
        #loop in each chart sheet and load in data frame for each chart sheet
        #convert and add chart of the correct chart_type
        for index,sheet in enumerate(chart_sheets):
            chart_df = self.chart_data_load(self.all_files.data_name,sheet)
            col_names = chart_df.columns.tolist()
            main_presentation,chart,col,row = self.create_chart(main_presentation,chart_df,col_names,chart_types[index],index)
            chart.Name = chart_sheets[index]
            num_of_cols.append(col)
            num_of_rows.append(row)
          
        #save the presentation in output_name - note this saves over any existing file
        main_presentation.save(self.all_files.output_name)
        #Cycle through all the charts in the presentation and apply the correct chart template from excel index sheet
        self.apply_chart_templates(chart_types,chart_labels,num_of_cols,num_of_rows)

        # Send done message when complete
        self.make_label("Done",1,1,150)



if __name__=="__main__":
    app = MainApp()
    app.mainloop()