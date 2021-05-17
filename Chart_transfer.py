# Chart transfer
# loads data from a specially formatted excel sheet
# and creates native Powerpoint versions in a new presentation
# A table in the "index" worksheet gives the sheet names for each set of chart data_name
# The "sheets" col has the sheet names and the "type" col has each chart type
# Other data can be added to this sheet over time to alter formatting of each chart

# Libaries for excel and PPT

import pandas as pd
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches

# UI for transfer
import tkinter as tk
from tkinter import *
from tkinter import filedialog
from os import path
import sys
sys.setrecursionlimit(20000)

from chart_type_data import chart_type_dict, chart_template_dict


#create UI menus
root = Tk()
root.title("Chart Transfer 0.1")
root.geometry('400x400')
menu = Menu(root)
root.config(menu=menu)

# load the config file with default loacations and the chart template files
import configparser

config = configparser.ConfigParser()
config.read("config.txt")
defaults = config['DEFAULTS']
chart_template_location = defaults['chart_template_location']
print(chart_template_location)
chart_templates = list(dict(config['CHART TEMPLATES']).values())

print(chart_templates)

file_name=""
output_name=defaults["output_name"]  #r"C:\Users\jamie\OneDrive - Market Prescience\python\powerpoint\output.pptx"
data_name=defaults["data_name"]  #"default_data.xlsx"
template_name=defaults["template_name"]  #"default_template.pptx"




# open file menu to select template file - just a plane pptx file rather than a pptm
def choose_template():
    global template_name
    template_name = filedialog.askopenfilename(title="Select a File",filetypes=(("Powerpoint","*.pptx"),("All types","*.*")))
    #my_label = Label(root, text=template_name).pack()
    label = tk.Label(text=template_name, bg="black", fg="white", wraplength = 150)
    label.grid(row=0, column=1, sticky="nsew")

# open file menu to select data file
def choose_data():
    global data_name
    data_name = filedialog.askopenfilename(title="Select a File",filetypes=(("Excel","*.xlsx"),("All types","*.*")))
    label = tk.Label(text=data_name, bg="black", fg="white", wraplength = 150)
    label.grid(row=1, column=1, sticky="nsew")

# open file menu to select output pointpoint file - will be written over
def choose_output():
    global output_name
    output_name = filedialog.askopenfilename(title="Select a File",filetypes=(("Powerpoint","*.pptx"),("All types","*.*")))
    label = tk.Label(text=output_name, bg="black", fg="white", wraplength = 150)
    label.grid(row=2, column=1, sticky="nsew")
    print(output_name)

def exit_prog():
    sys.exit()


# tests all whether there is a template, a data file and a output destination

def transfer_all():
    global output_name, data_name, template_name
    if output_name=="" or data_name =="" or template_name =="":
        label = tk.Label(text="Please select all files", bg="black", fg="white")
        label.grid(row=3, column=1, sticky="nsew")
    else:
        execute_transfer()

def open_chart_index(file_name):
    #open the excel file and grab the chart index infomation from index worksheet
    #this data contains the excel sheet names for each chart and chart types
    df = pd.read_excel (file_name, sheet_name="index",engine='openpyxl')
    return df


def open_presentation(template):
    #open the presentation template - this will be user defined
    prs = Presentation(template)
    return prs

def chart_data_load(file_name,sheet):
    # Open excel file and load data from sheet into pandas data frame
    df = pd.read_excel (file_name, sheet_name=sheet,engine='openpyxl')
    return df



def create_chart(prs,df,col_names,chart_type):
    # create the slide and add the chart data
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    chart_data = CategoryChartData()
    chart_data.categories = (df['category'].values.tolist())
    # take the category column off the chart data list
    col_names.pop(0)

    #count col and row Numbers
    num_of_cols = len(col_names)
    num_of_rows = len(df['category'].values.tolist())


    # loop through the columns of data and add to the active chart_data note that the data needs to be converted from list to tuple
    for x in range(num_of_cols):
        chart_data.add_series(str(col_names[x]),tuple(df[col_names[x]].values.tolist()))

    # add chart to slide --------------------
    x, y, cx, cy = Inches(1), Inches(2), Inches(10.5), Inches(5)
    # lookup chart_type info from dictionary of excel chart types
    obj = chart_type_dict[chart_type]
    chart=slide.shapes.add_chart(
        obj, x, y, cx, cy, chart_data
        ).chart
    slide.shapes.Name = "chart1"
    return prs,chart,num_of_cols,num_of_rows

def format_chart(chart,chart_label):
    # create the slide and add the chart data
    from pptx.enum.chart import XL_LABEL_POSITION
    plot = chart.plots[0]
    plot.has_data_labels = True
    data_labels = plot.data_labels

    if chart_label=="per":
        data_labels.number_format = "%"
    else:
        data_labels.number_format ='0'

    return chart

def apply_chart_templates(types,labels,num_of_cols,num_of_rows):
    # Uses win32com to open powerpoint and apply chart templates to all the charts in the output file
    import win32com.client
    from pathlib import Path
    pptApp = win32com.client.Dispatch ("Powerpoint.Application")
    pptApp.Visible = True


    file = Path(output_name)
    prs = pptApp.Presentations.Open(file, ReadOnly = False)

    # go through all the slides in the Presentation
    # examine each of the "Shapes" and find all the charts - then apply the template
    for index,current_slide in enumerate(prs.Slides):
        for shp in current_slide.Shapes:
            if shp.Type == 3:   #type 3 is msochart
                x=chart_template_dict[types[index-1]]
                print(types[index-1],chart_templates[x-1])
                shp.Chart.ApplyChartTemplate(Path(chart_template_location+chart_templates[x-1]))
                print(index)
                shp.Chart.Axes(2).MaximumScaleIsAuto = True # 2 - is the xlValue constant number
                shp.Chart.Axes(2).MinimumScaleIsAuto = True
                #shp.Chart.Axes(2).TickLabels.Font.ColorIndex = 3
                for col in range(1,num_of_cols[index-1]+1):
                    plot = shp.Chart.SeriesCollection(col)
                    plot.HasDataLabels = True
                    for row in range(1,num_of_rows[index-1]+1):
                        if labels[index-1]=="per":
                            plot.DataLabels(row).NumberFormat = "0%"
                        else:
                            plot.DataLabels(row).NumberFormat = "0"



def execute_transfer():
    #load the chart index dataframe
    chart_index_df = open_chart_index(data_name)

    #open the template presentation
    main_presentation = open_presentation(template_name)

    #create all the charts from the chart index
    chart_sheets = chart_index_df['sheets'].values.tolist() # put the sheet names into list
    chart_types = chart_index_df['type'].values.tolist() #put chart types into list
    chart_labels = chart_index_df['label'].values.tolist() # put the chart label type into list (numbers or counts are num and % are per)
    num_of_rows = []
    num_of_cols = []

    #loop in each chart sheet and load in data frame for each chart sheet
    #convert and add chart of the correct chart_type
    for index,sheet in enumerate(chart_sheets):
        chart_df = chart_data_load(data_name,sheet)
        col_names = chart_df.columns.tolist()
        main_presentation,chart,col,row = create_chart(main_presentation,chart_df,col_names,chart_types[index])
        chart.Name = chart_sheets[index]
        num_of_cols.append(col)
        num_of_rows.append(row)
        #format_chart(chart,chart_labels[index])

    print(num_of_cols,num_of_rows)
    #save the presentation in output_name - note this saves over any existing file
    main_presentation.save(output_name)
    label = tk.Label(text="Done", bg="red", fg="white")
    label.grid(row=3, column=1, sticky="nsew")
    apply_chart_templates(chart_types,chart_labels,num_of_cols,num_of_rows)



    sys.exit()

#create UI col and rows

root.rowconfigure([0,1,2,3], minsize=100)
root.columnconfigure([0, 1], minsize=200)

button1 = tk.Button(text="Template:", bg="black", fg="white",command=choose_template)
button2 = tk.Button(text="Data:", bg="black", fg="white", command=choose_data)
button3 = tk.Button(text="Output", bg="black", fg="white", command = choose_output)
button4 = tk.Button(text="Execute", bg="black", fg="white", command = transfer_all)

button1.grid(row=0, column=0, sticky="nsew")
button2.grid(row=1, column=0, sticky="nsew")
button3.grid(row=2, column=0, sticky="nsew")
button4.grid(row=3, column=0, sticky="nsew")

label = tk.Label(text=template_name, bg="black", fg="white", wraplength = 150)
label.grid(row=0, column=1, sticky="nsew")

label = tk.Label(text=data_name, bg="black", fg="white", wraplength = 150)
label.grid(row=1, column=1, sticky="nsew")

label = tk.Label(text=output_name, bg="black", fg="white", wraplength = 150)
label.grid(row=2, column=1, sticky="nsew")

# create file menu
file_menu = Menu(menu)
menu.add_cascade(label='File', menu=file_menu)
#file_menu.add_command(label='Options', command=options)
file_menu.add_command(label='Exit',command=exit_prog)

# create run menu
run_menu = Menu(menu)
menu.add_cascade(label='Run', menu=run_menu)
run_menu.add_command(label='Transfer', command=transfer_all)


clicked = IntVar()
clicked.set(4)


root.mainloop()
